"""Generate the EduVoice hackathon pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x0A, 0x0A, 0x1A)
PURPLE = RGBColor(0x81, 0x8C, 0xF8)
LIGHT_PURPLE = RGBColor(0xA5, 0xB4, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED_ACCENT = RGBColor(0xF8, 0x71, 0x71)
CARD_BG = RGBColor(0x16, 0x16, 0x2E)

W = Inches(13.333)
H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=GRAY, bullet_color=PURPLE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def slide_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 subtitle, font_size=18, color=GRAY)


# ── Build presentation ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Slide 1: Title ──────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(sl, DARK_BG)
add_text(sl, Inches(0), Inches(1.5), W, Inches(1),
         "EduVoice", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(2.6), W, Inches(0.8),
         "Voice-First Learning Assistant", font_size=28, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(3.8), W, Inches(0.6),
         "Ask anything. Learn by listening. In any language.",
         font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(5.5), W, Inches(0.5),
         "Voice AI Hackathon  |  Track 3: Accessibility & Societal Impact",
         font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(6.2), W, Inches(0.5),
         "Built by Suvarna Sanapathi", font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 2: Problem ────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
slide_title(sl, "The Problem")
add_rounded_rect(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5))
items = [
    "Millions of students lack access to quality tutors, especially in rural areas",
    "Text-heavy apps exclude low-literacy and visually impaired learners",
    "Language barriers — most ed-tech is English-only",
    "Students need instant, judgment-free help at any hour",
    "Existing solutions require screens, typing, and expensive subscriptions",
]
add_bullet_list(sl, Inches(1.2), Inches(2.3), Inches(10.5), Inches(4.0),
                [f"  {item}" for item in items], font_size=18, color=GRAY)

# ── Slide 3: Solution ───────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
slide_title(sl, "Our Solution", "EduVoice — learn by just speaking")

cards = [
    ("Voice-First", "No typing, no screens needed.\nJust speak and learn."),
    ("RAG-Powered", "Retrieves real knowledge\nbefore answering — accurate."),
    ("Multilingual", "Speak in Hindi, English,\nor any language."),
    ("Always Available", "24/7 patient tutor that\nnever judges."),
]
for i, (title, desc) in enumerate(cards):
    left = Inches(0.8 + i * 3.1)
    add_rounded_rect(sl, left, Inches(2.2), Inches(2.8), Inches(3.5))
    add_text(sl, left, Inches(2.5), Inches(2.8), Inches(0.5),
             title, font_size=20, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(sl, left + Inches(0.2), Inches(3.2), Inches(2.4), Inches(2.0),
             desc, font_size=15, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 4: Architecture ───────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
slide_title(sl, "Architecture", "End-to-end voice RAG pipeline")

steps = [
    ("Student\nSpeaks", PURPLE),
    ("Vapi\n(STT)", RGBColor(0x60, 0xA5, 0xFA)),
    ("FastAPI\nServer", GREEN),
    ("Qdrant\n(Search)", RGBColor(0xFB, 0xBF, 0x24)),
    ("OpenAI\n(LLM)", RED_ACCENT),
    ("Vapi\n(TTS)", RGBColor(0x60, 0xA5, 0xFA)),
    ("Student\nHears", PURPLE),
]
for i, (label, color) in enumerate(steps):
    left = Inches(0.6 + i * 1.75)
    top = Inches(3.0)
    add_rounded_rect(sl, left, top, Inches(1.4), Inches(1.6))
    add_text(sl, left, top + Inches(0.3), Inches(1.4), Inches(1.0),
             label, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(sl, left + Inches(1.4), top + Inches(0.4), Inches(0.4), Inches(0.5),
                 "\u2192", font_size=24, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_text(sl, Inches(0.8), Inches(5.4), Inches(11), Inches(0.5),
         "Speech \u2192 Text \u2192 Embed \u2192 Vector Search \u2192 LLM Generation \u2192 Speech",
         font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 5: Tech Stack ─────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
slide_title(sl, "Tech Stack")

techs = [
    ("Vapi", "Voice AI platform\nSTT + TTS + Call management"),
    ("Qdrant", "Vector database\nSemantic knowledge retrieval"),
    ("OpenAI", "Embeddings (text-embedding-3-small)\nLLM (gpt-4o-mini)"),
    ("FastAPI", "Python backend\nWebhook + Custom LLM endpoint"),
]
for i, (name, desc) in enumerate(techs):
    left = Inches(0.8 + (i % 2) * 6.0)
    top = Inches(2.2 + (i // 2) * 2.2)
    add_rounded_rect(sl, left, top, Inches(5.5), Inches(1.8))
    add_text(sl, left + Inches(0.3), top + Inches(0.2), Inches(4.8), Inches(0.5),
             name, font_size=22, color=PURPLE, bold=True)
    add_text(sl, left + Inches(0.3), top + Inches(0.8), Inches(4.8), Inches(0.8),
             desc, font_size=14, color=GRAY)

# ── Slide 6: Features ───────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
slide_title(sl, "Key Features")

features = [
    "RAG-powered answers — retrieves knowledge before answering, not just LLM guessing",
    "Multi-subject — Physics, Chemistry, Math, Biology, History, CS, and more",
    "Conversation memory — maintains context across turns within a call",
    "Multilingual — responds in the same language the student speaks",
    "Voice-optimized — concise, spoken-friendly responses",
    "Encouraging tone — patient and supportive, perfect for learners",
    "Web UI — one-click browser call, no app install needed",
]
add_bullet_list(sl, Inches(1.2), Inches(2.0), Inches(10.5), Inches(5.0),
                [f"  {f}" for f in features], font_size=17, color=GRAY)

# ── Slide 7: Impact ─────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
slide_title(sl, "Societal Impact", "Making education accessible through voice")

impacts = [
    ("Rural Students", "No need for expensive devices or fast internet — just a phone call"),
    ("Low-Literacy Learners", "Voice removes the barrier of reading and typing"),
    ("Visually Impaired", "Fully accessible through audio — no screen dependency"),
    ("Multilingual Communities", "Learn in your own language, not just English"),
]
for i, (title, desc) in enumerate(impacts):
    top = Inches(2.2 + i * 1.15)
    add_rounded_rect(sl, Inches(0.8), top, Inches(11.5), Inches(1.0))
    add_text(sl, Inches(1.2), top + Inches(0.1), Inches(3.5), Inches(0.7),
             title, font_size=18, color=PURPLE, bold=True)
    add_text(sl, Inches(4.5), top + Inches(0.1), Inches(7.5), Inches(0.7),
             desc, font_size=16, color=GRAY)

# ── Slide 8: Demo ───────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
add_text(sl, Inches(0), Inches(2.0), W, Inches(1),
         "Live Demo", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(3.3), W, Inches(0.6),
         '"What is Newton\'s third law?"', font_size=24, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(4.0), W, Inches(0.6),
         '"Explain photosynthesis in Hindi"', font_size=24, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(4.7), W, Inches(0.6),
         '"What is DNA?"', font_size=24, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(5.8), W, Inches(0.5),
         "Try it: click Talk on the web UI", font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 9: Thank You ──────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DARK_BG)
add_text(sl, Inches(0), Inches(2.0), W, Inches(1),
         "Thank You!", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(3.3), W, Inches(0.6),
         "EduVoice — Making learning accessible, one voice at a time.",
         font_size=20, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(4.8), W, Inches(0.4),
         "github.com/suvvvv/EduVoice", font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text(sl, Inches(0), Inches(5.3), W, Inches(0.4),
         "Suvarna Sanapathi  |  suvarnasanapathi2001@gmail.com",
         font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────
out = "/Users/latha/Public/my-projects/EduVoice/EduVoice_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
